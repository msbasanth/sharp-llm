from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

pptx_path = r'paper/GenAI_SHARP-LLM_A_Framework_for_Healthcare_Vulnerabilitty_Analysis_AM.SC.R4CSE25007-PT_v2.pptx'
prs = Presentation(pptx_path)

print(f'Total slides: {len(prs.slides)}')
print()

for i, slide in enumerate(prs.slides, 1):
    layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'
    print(f'--- Slide {i} (layout: {layout_name}) ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    # Show font size if available
                    size_info = ""
                    for run in para.runs:
                        if run.font.size:
                            size_info = f" [{run.font.size.pt:.0f}pt]"
                            break
                    if len(text) > 120:
                        text = text[:120] + '...'
                    print(f'  {text}{size_info}')
        elif shape.has_table:
            tbl = shape.table
            print(f'  [TABLE: {len(tbl.rows)} rows x {len(tbl.columns)} cols]')
        elif shape.shape_type == 13:
            print(f'  [IMAGE]')
    print()

for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    # A divider slide has exactly 2 text items: a number (01-07) and a section title
    if len(texts) == 2 and texts[0].isdigit() and len(texts[0]) <= 2:
        divider_data[i] = (texts[0], texts[1])
        print(f"  Divider at slide {i+1}: '{texts[0]}' - '{texts[1]}'")

print(f"\nFound {len(divider_data)} divider slides")

# Step 2: For each divider, find the title shape on the next slide
# Use multiple heuristics: placeholder type, then topmost short text shape
def find_title_shape(slide):
    """Find the title shape on a slide using placeholders or position."""
    # Try 1: Look for a placeholder with idx 0 or 1 (standard title placeholders)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx in (0, 1):
            if shape.has_text_frame and shape.text_frame.text.strip():
                return shape
    
    # Try 2: Find topmost text shape that has short text (likely a title, not body content)
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            # Get first non-empty paragraph text
            first_text = ""
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    first_text = para.text.strip()
                    break
            # Title candidates: short text, no bullet markers
            if first_text and not first_text.startswith('▸') and not first_text.startswith('◦'):
                text_shapes.append((shape.top, shape, first_text))
    
    if text_shapes:
        # Sort by vertical position (topmost first)
        text_shapes.sort(key=lambda x: x[0])
        # Return topmost non-bullet text shape
        return text_shapes[0][1]
    
    return None

# Step 3: Merge divider text into next slide's title
for div_idx in sorted(divider_data.keys()):
    number, section_title = divider_data[div_idx]
    next_idx = div_idx + 1
    if next_idx >= len(prs.slides):
        continue
    
    next_slide = prs.slides[next_idx]
    title_shape = find_title_shape(next_slide)
    
    if title_shape is None:
        print(f"  WARNING: No title shape found on slide {next_idx+1}")
        continue
    
    tf = title_shape.text_frame
    
    # Get existing title text (first non-empty paragraph)
    existing_title = ""
    for para in tf.paragraphs:
        if para.text.strip():
            existing_title = para.text.strip()
            break
    
    # Get existing formatting from runs
    existing_font_name = None
    existing_font_bold = None
    existing_font_color = None
    existing_font_size = None
    for para in tf.paragraphs:
        for run in para.runs:
            if run.text.strip():
                existing_font_name = run.font.name
                existing_font_bold = run.font.bold
                existing_font_size = run.font.size
                try:
                    existing_font_color = run.font.color.rgb
                except (AttributeError, TypeError):
                    existing_font_color = None
                break
        if existing_font_name:
            break
    
    # Clear and rebuild the text frame with two paragraphs
    tf.clear()
    
    # Paragraph 1: Section title from divider (original font size)
    p1 = tf.paragraphs[0]
    run1 = p1.add_run()
    run1.text = section_title
    if existing_font_size:
        run1.font.size = existing_font_size
    if existing_font_bold is not None:
        run1.font.bold = existing_font_bold
    if existing_font_color:
        run1.font.color.rgb = existing_font_color
    if existing_font_name:
        run1.font.name = existing_font_name
    
    # Paragraph 2: Original title at 24pt
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = existing_title
    run2.font.size = Pt(24)
    if existing_font_bold is not None:
        run2.font.bold = existing_font_bold
    if existing_font_color:
        run2.font.color.rgb = existing_font_color
    if existing_font_name:
        run2.font.name = existing_font_name
    
    print(f"  Slide {next_idx+1}: '{section_title}' + '{existing_title}' (24pt)")

# Step 4: Remove divider slides (reverse order to preserve indices)
for div_idx in sorted(divider_data.keys(), reverse=True):
    rId = prs.slides._sldIdLst[div_idx].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[div_idx]
    print(f"  Removed divider (was slide {div_idx+1})")

# Save
out_path = r'paper/GenAI_SHARP-LLM_A_Framework_for_Healthcare_Vulnerabilitty_Analysis_AM.SC.R4CSE25007-PT_v2.pptx'
prs.save(out_path)
print(f"\nDone! Saved: {out_path}")
print(f"Slides: 30 → {len(prs.slides)}")
