from pptx import Presentation
from pptx.util import Pt

pptx_path = r'paper/GenAI_SHARP-LLM_A_Framework_for_Healthcare_Vulnerabilitty_Analysis_AM.SC.R4CSE25007-PT_v2.pptx'
prs = Presentation(pptx_path)

# Search for the slide with "Gap 1:" content
slide_idx = None
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "Gap 1:" in shape.text_frame.text:
                slide_idx = i
                break
    if slide_idx is not None:
        break

if slide_idx is None:
    print("ERROR: Research Gaps slide not found!")
else:
    slide = prs.slides[slide_idx]
    print(f"=== Slide {slide_idx+1} — Research Gaps ===\n")

for shape_idx, shape in enumerate(slide.shapes):
    print(f"Shape {shape_idx}: type={shape.shape_type}, name='{shape.name}'")
    print(f"  Position: top={shape.top}, left={shape.left}, width={shape.width}, height={shape.height}")
    
    if shape.has_text_frame:
        tf = shape.text_frame
        print(f"  Text frame: {len(tf.paragraphs)} paragraphs")
        for p_idx, para in enumerate(tf.paragraphs):
            full_text = para.text
            if full_text.strip():
                size_info = ""
                for run in para.runs:
                    if run.font.size:
                        size_info = f" [{run.font.size.pt:.0f}pt]"
                        break
                print(f"    P{p_idx}: '{full_text}'{size_info}")
                # Show individual runs
                for r_idx, run in enumerate(para.runs):
                    rs = run.font.size.pt if run.font.size else "inherited"
                    print(f"      Run{r_idx}: '{run.text}' size={rs} bold={run.font.bold}")
    elif shape.has_table:
        print(f"  [TABLE]")
    elif shape.shape_type == 13:
        print(f"  [IMAGE]")
    print()
