from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

pptx_path = r'paper/GenAI_SHARP-LLM_A_Framework_for_Healthcare_Vulnerabilitty_Analysis_AM.SC.R4CSE25007-PT_v2.pptx'
prs = Presentation(pptx_path)

print(f'Total slides: {len(prs.slides)}')
print()

for i, slide in enumerate(prs.slides, 1):
    layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'
    print(f'--- Slide {i} (layout: {layout_name}) ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    size_info = ""
                    for run in para.runs:
                        if run.font.size:
                            size_info = f" [{run.font.size.pt:.0f}pt]"
                            break
                    if len(text) > 120:
                        text = text[:120] + '...'
                    print(f'  {text}{size_info}')
        elif shape.has_table:
            tbl = shape.table
            print(f'  [TABLE: {len(tbl.rows)} rows x {len(tbl.columns)} cols]')
        elif shape.shape_type == 13:
            print(f'  [IMAGE]')
    print()
